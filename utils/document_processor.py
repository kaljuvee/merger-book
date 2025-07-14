"""
Document processing utilities for Merger Book MVP
Handles PDF, Word, and PowerPoint document parsing and content extraction
"""

import os
import tempfile
from typing import Dict, List, Optional, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import streamlit as st
from PIL import Image
import io
import base64

class DocumentProcessor:
    """Handles document parsing and content extraction"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_docx,  # Will attempt to process as docx
            'pptx': self._process_pptx,
            'ppt': self._process_pptx   # Will attempt to process as pptx
        }
    
    def process_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Process a document and extract content
        
        Args:
            file_path: Path to the document file
            file_type: Type of document (pdf, docx, pptx, etc.)
        
        Returns:
            Dictionary containing extracted content and metadata
        """
        try:
            if file_type.lower() not in self.supported_formats:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            processor = self.supported_formats[file_type.lower()]
            result = processor(file_path)
            
            # Add common metadata
            result.update({
                'file_path': file_path,
                'file_type': file_type,
                'processing_status': 'completed',
                'error_message': None
            })
            
            return result
            
        except Exception as e:
            return {
                'file_path': file_path,
                'file_type': file_type,
                'processing_status': 'error',
                'error_message': str(e),
                'text_content': '',
                'page_count': 0,
                'metadata': {}
            }
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF document"""
        text_content = []
        metadata = {}
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Extract metadata
            if pdf_reader.metadata:
                metadata = {
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'producer': pdf_reader.metadata.get('/Producer', ''),
                    'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')),
                    'modification_date': str(pdf_reader.metadata.get('/ModDate', ''))
                }
            
            # Extract text from each page
            page_count = len(pdf_reader.pages)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append({
                            'page_number': page_num + 1,
                            'content': page_text.strip()
                        })
                except Exception as e:
                    st.warning(f"Could not extract text from page {page_num + 1}: {str(e)}")
        
        return {
            'text_content': '\n\n'.join([page['content'] for page in text_content]),
            'pages': text_content,
            'page_count': page_count,
            'metadata': metadata,
            'document_type': 'pdf'
        }
    
    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process Word document"""
        try:
            doc = Document(file_path)
            
            # Extract text content
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract metadata
            metadata = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'keywords': doc.core_properties.keywords or '',
                'comments': doc.core_properties.comments or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
                'last_modified_by': doc.core_properties.last_modified_by or ''
            }
            
            # Extract tables if any
            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_content.append(table_data)
            
            return {
                'text_content': '\n\n'.join(text_content),
                'paragraphs': text_content,
                'tables': tables_content,
                'page_count': 1,  # Word docs don't have clear page breaks in python-docx
                'metadata': metadata,
                'document_type': 'docx'
            }
            
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            
            # Extract text content from slides
            slides_content = []
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slide_content = {
                    'slide_number': slide_num + 1,
                    'content': '\n'.join(slide_text)
                }
                
                slides_content.append(slide_content)
                all_text.extend(slide_text)
            
            # Extract metadata
            metadata = {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'subject': prs.core_properties.subject or '',
                'keywords': prs.core_properties.keywords or '',
                'comments': prs.core_properties.comments or '',
                'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                'modified': str(prs.core_properties.modified) if prs.core_properties.modified else '',
                'last_modified_by': prs.core_properties.last_modified_by or '',
                'slide_count': len(prs.slides)
            }
            
            return {
                'text_content': '\n\n'.join(all_text),
                'slides': slides_content,
                'page_count': len(prs.slides),
                'metadata': metadata,
                'document_type': 'pptx'
            }
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint presentation: {str(e)}")
    
    def extract_business_features(self, text_content: str) -> Dict[str, Any]:
        """
        Extract basic business features from text content
        This is a simple rule-based extraction that will be enhanced with AI
        """
        features = {
            'revenue_mentions': [],
            'industry_keywords': [],
            'geographic_mentions': [],
            'employee_mentions': [],
            'technology_mentions': [],
            'financial_metrics': []
        }
        
        # Simple keyword-based extraction
        text_lower = text_content.lower()
        
        # Revenue patterns
        revenue_keywords = ['revenue', 'sales', 'income', 'earnings', 'turnover']
        for keyword in revenue_keywords:
            if keyword in text_lower:
                features['revenue_mentions'].append(keyword)
        
        # Industry keywords
        industry_keywords = [
            'technology', 'fintech', 'healthcare', 'manufacturing', 'retail',
            'e-commerce', 'software', 'saas', 'artificial intelligence', 'ai',
            'machine learning', 'blockchain', 'cryptocurrency', 'biotech'
        ]
        for keyword in industry_keywords:
            if keyword in text_lower:
                features['industry_keywords'].append(keyword)
        
        # Geographic mentions
        geo_keywords = [
            'united states', 'usa', 'europe', 'asia', 'global', 'international',
            'north america', 'california', 'new york', 'london', 'singapore'
        ]
        for keyword in geo_keywords:
            if keyword in text_lower:
                features['geographic_mentions'].append(keyword)
        
        # Employee mentions
        employee_keywords = ['employees', 'staff', 'team', 'workforce', 'personnel']
        for keyword in employee_keywords:
            if keyword in text_lower:
                features['employee_mentions'].append(keyword)
        
        # Technology mentions
        tech_keywords = [
            'platform', 'api', 'cloud', 'mobile', 'web', 'database',
            'analytics', 'automation', 'digital', 'innovation'
        ]
        for keyword in tech_keywords:
            if keyword in text_lower:
                features['technology_mentions'].append(keyword)
        
        return features
    
    def save_uploaded_file(self, uploaded_file, upload_folder: str) -> str:
        """
        Save uploaded file to disk and return file path
        
        Args:
            uploaded_file: Streamlit uploaded file object
            upload_folder: Directory to save the file
        
        Returns:
            Path to saved file
        """
        # Ensure upload folder exists
        os.makedirs(upload_folder, exist_ok=True)
        
        # Generate unique filename
        file_extension = uploaded_file.name.split('.')[-1].lower()
        timestamp = int(time.time())
        filename = f"{timestamp}_{uploaded_file.name}"
        file_path = os.path.join(upload_folder, filename)
        
        # Save file
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        return file_path
    
    def get_file_info(self, uploaded_file) -> Dict[str, Any]:
        """Get information about uploaded file"""
        return {
            'filename': uploaded_file.name,
            'file_type': uploaded_file.name.split('.')[-1].lower(),
            'file_size': uploaded_file.size,
            'mime_type': uploaded_file.type
        }

# Import time for timestamp generation
import time

